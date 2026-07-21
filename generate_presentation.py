import sys
import os
import subprocess

# Auto-install python-pptx if missing
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx not found. Installing now...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set to widescreen 16:9 format
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6] # Blank layout for custom styling

    # Theme Colors
    BG_COLOR = RGBColor(5, 5, 8)       # #050508 Deep Dark
    TEXT_WHITE = RGBColor(244, 244, 249) # #f4f4f9 Off White
    CYAN = RGBColor(6, 182, 212)       # #06b6d4 Cyber Cyan
    INDIGO = RGBColor(99, 102, 241)     # #6366f1 Indigo
    MUTED = RGBColor(138, 153, 173)     # #8a99ad Grey
    CARD_BG = RGBColor(15, 15, 26)      # #0f0f1a Secondary Dark

    def apply_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_title(slide, text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(1))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        return title_box

    # --- SLIDE 1: TITLE SLIDE ---
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)

    # Accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(3.5), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.color.rgb = CYAN

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(2.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TwinFusion AI"
    p.font.name = 'Arial'
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Co-Piloting the Campus Transition with a Personal AI Digital Twin"
    p2.font.name = 'Arial'
    p2.font.size = Pt(22)
    p2.font.color.rgb = CYAN
    p2.space_before = Pt(14)

    info_box = slide.shapes.add_textbox(Inches(0.75), Inches(5.8), Inches(6), Inches(1))
    tf_info = info_box.text_frame
    p_team = tf_info.paragraphs[0]
    p_team.text = "Google Student Ambassador Hackathon"
    p_team.font.name = 'Arial'
    p_team.font.size = Pt(14)
    p_team.font.color.rgb = MUTED
    p_names = tf_info.add_paragraph()
    p_names.text = "Pitch Presentation Deck"
    p_names.font.name = 'Arial'
    p_names.font.size = Pt(12)
    p_names.font.color.rgb = MUTED

    # --- SLIDE 2: THE PROBLEM ---
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_title(slide, "The Fresher's Dilemma")

    problems = [
        ("Information Overload", "Important alerts and announcements are scattered across chaos channels (WhatsApp, print boards, messy sites)."),
        ("Broken Campus Navigation", "Lost classrooms, faculty offices, and admin departments cause stress and delays during early weeks."),
        ("Missed Opportunities", "Students miss crucial club entries, workshops, and hackathons simply due to lack of visibility.")
    ]

    for i, (title, desc) in enumerate(problems):
        x = Inches(0.75 + i * 4.0)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(3.7), Inches(4.2))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = INDIGO if i == 0 else MUTED
        card.line.width = Pt(1.5)

        tb = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.2), Inches(3.3), Inches(3.8))
        tf = tb.text_frame
        tf.word_wrap = True
        
        # Number indicator
        p0 = tf.paragraphs[0]
        p0.text = f"0{i+1}"
        p0.font.name = 'Courier New'
        p0.font.size = Pt(28)
        p0.font.bold = True
        p0.font.color.rgb = CYAN if i == 1 else INDIGO

        p1 = tf.add_paragraph()
        p1.text = title
        p1.font.name = 'Arial'
        p1.font.size = Pt(20)
        p1.font.bold = True
        p1.font.color.rgb = TEXT_WHITE
        p1.space_before = Pt(14)
        p1.space_after = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(14)
        p2.font.color.rgb = MUTED
        p2.line_spacing = 1.3

    # --- SLIDE 3: THE SOLUTION ---
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_title(slide, "Our Solution: The AI Digital Twin")

    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.5), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Instead of presenting the same static information to everyone, TwinFusion AI models a personal digital companion for each student."
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.color.rgb = TEXT_WHITE
    p.line_spacing = 1.3
    p.space_after = Pt(20)

    points = [
        ("9-Question Profile Setup", "Twin calibrates based on department, coding experience, and aspirations."),
        ("Timetable-Centric View", "Your planner is built dynamically matching your daily course slots."),
        ("Proactive AI Companion", "Acts like a smart senior in your pocket, predicting your queries.")
    ]
    for pt_t, pt_d in points:
        p_t = tf.add_paragraph()
        p_t.text = f"• {pt_t}: "
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(16)
        p_t.font.bold = True
        p_t.font.color.rgb = CYAN
        p_t.space_before = Pt(10)
        
        p_d = tf.add_paragraph()
        p_d.text = f"  {pt_d}"
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = MUTED

    # Graphic box representing the constellation link
    g_card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(1.8), Inches(4.78), Inches(4.5))
    g_card.fill.solid()
    g_card.fill.fore_color.rgb = CARD_BG
    g_card.line.color.rgb = CYAN
    
    g_tb = slide.shapes.add_textbox(Inches(7.9), Inches(2.0), Inches(4.58), Inches(4.1))
    g_tf = g_tb.text_frame
    g_tf.word_wrap = True
    g_p = g_tf.paragraphs[0]
    g_p.text = "THE CONSTELLATION SYNC"
    g_p.font.name = 'Arial'
    g_p.font.size = Pt(14)
    g_p.font.bold = True
    g_p.font.color.rgb = CYAN
    g_p.alignment = PP_ALIGN.CENTER
    g_p.space_after = Pt(24)

    visuals = ["YOU Node  <==== Pulse Telemetry ====>  TWIN Node",
               "🛰️ Interests Node synced",
               "🛰️ Skills calibrated",
               "🛰️ Academic Schedule mapped",
               "🛰️ Targets cataloged"]
    for vis in visuals:
        g_v = g_tf.add_paragraph()
        g_v.text = vis
        g_v.font.name = 'Courier New'
        g_v.font.size = Pt(13)
        g_v.font.color.rgb = TEXT_WHITE
        g_v.alignment = PP_ALIGN.CENTER
        g_v.space_before = Pt(14)

    # --- SLIDE 4: OPPORTUNITY RADAR ---
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_title(slide, "Key Feature: Opportunity Radar")

    tb = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Smart Matchmaking & Filtering"
    p.font.name = 'Arial'
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(14)

    bullets = [
        ("Interests Mapping", "Checks coding preferences (e.g. Python) to alert about hackathons, ignoring noise."),
        ("Deadlines Tracking", "Reminds you days ahead about club entries and events before they close."),
        ("Google Ambassadorship Radar", "Directly targets Google-sponsored hackathons and local coding challenges.")
    ]
    for b_t, b_d in bullets:
        p_bt = tf.add_paragraph()
        p_bt.text = f"✔ {b_t}"
        p_bt.font.name = 'Arial'
        p_bt.font.size = Pt(16)
        p_bt.font.bold = True
        p_bt.font.color.rgb = CYAN
        p_bt.space_before = Pt(12)
        
        p_bd = tf.add_paragraph()
        p_bd.text = b_d
        p_bd.font.name = 'Arial'
        p_bd.font.size = Pt(14)
        p_bd.font.color.rgb = MUTED

    # Calibration slider mock visual
    m_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.75), Inches(1.8), Inches(5.83), Inches(4.5))
    m_card.fill.solid()
    m_card.fill.fore_color.rgb = CARD_BG
    m_card.line.color.rgb = INDIGO

    m_tb = slide.shapes.add_textbox(Inches(6.95), Inches(2.0), Inches(5.43), Inches(4.1))
    m_tf = m_tb.text_frame
    m_tf.word_wrap = True
    
    mp = m_tf.paragraphs[0]
    mp.text = "Interactive Dashboard Preview"
    mp.font.name = 'Arial'
    mp.font.size = Pt(16)
    mp.font.bold = True
    mp.font.color.rgb = TEXT_WHITE
    mp.space_after = Pt(20)

    stages = [
        ("0% - 29% Calibration", "Awaiting Profile Sync..."),
        ("30% - 64% Syncing", "Learning Interests: Python, AI detected"),
        ("65% - 89% Matches Active", "Deadlines sync: AI Club registers open!"),
        ("90% - 100% Fully Mapped", "Google GenAI Hackathon (98% skill match!)")
    ]
    for st, sd in stages:
        s_p = m_tf.add_paragraph()
        s_p.text = f"{st}: {sd}"
        s_p.font.name = 'Arial'
        s_p.font.size = Pt(13)
        s_p.font.color.rgb = CYAN if "90%" in st else MUTED
        s_p.space_before = Pt(8)

    # --- SLIDE 5: TECHNICAL ARCHITECTURE ---
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)
    add_title(slide, "Technical Architecture & Google Tech Stack")

    blocks = [
        ("Frontend UI", "React JS & Vite\nCSS Grid Transforms\nMouse 3D Parallax Canvas", CYAN),
        ("Cognitive Layer", "Google Gemini API\nContext processing\nCustom senior advisor responses", INDIGO),
        ("Backend System", "Node.js & Express\nREST API integrations\nDatabase coordination", CYAN),
        ("Database Storage", "MySQL Database\nUser credentials\nSync calendar schedules", MUTED)
    ]

    for i, (title, desc, color) in enumerate(blocks):
        x = Inches(0.75 + i * 2.95)
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.3), Inches(2.75), Inches(3.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = color
        card.line.width = Pt(2)

        tb = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.5), Inches(2.55), Inches(3.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = 'Arial'
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        p.space_after = Pt(14)

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.name = 'Arial'
        p2.font.size = Pt(13)
        p2.font.color.rgb = MUTED
        p2.line_spacing = 1.3

    # --- SLIDE 6: CONCLUSION ---
    slide = prs.slides.add_slide(blank_layout)
    apply_background(slide)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(2.2), Inches(3.5), Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CYAN
    shape.line.color.rgb = CYAN

    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "TwinFusion AI"
    p.font.name = 'Arial'
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = TEXT_WHITE

    p2 = tf.add_paragraph()
    p2.text = "Your digital senior. Your campus journey, co-piloted."
    p2.font.name = 'Arial'
    p2.font.size = Pt(24)
    p2.font.color.rgb = CYAN
    p2.space_before = Pt(18)

    p3 = tf.add_paragraph()
    p3.text = "Let's make campus onboarding simple, one digital twin at a time."
    p3.font.name = 'Arial'
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED
    p3.space_before = Pt(12)

    desktop_path = os.path.expanduser("~/Desktop")
    output_path = os.path.join(desktop_path, "TwinFusion_AI_Pitch.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_deck()
